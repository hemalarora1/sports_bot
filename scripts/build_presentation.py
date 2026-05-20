"""Build the CS225A PickleBot progress deck.

Generates `sports_bot/CS225A PickleBot Progress v2.pptx` from python-pptx
primitives. Numbers and behavior reflect the current state of
`sports_bot/state_machine/` and `sports_bot/optitrack/`. Re-run after any
tuning change to keep the deck honest.

Run:
    /opt/homebrew/Caskroom/miniconda/base/envs/opensai/bin/python \
        sports_bot/scripts/build_presentation.py
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt


# ----- paths -----------------------------------------------------------------

REPO_ROOT = Path(__file__).resolve().parents[1]
MEDIA_DIR = REPO_ROOT / "media"
OUTPUT = REPO_ROOT / "CS225A PickleBot Progress v2.pptx"

MOUNT_ISO = MEDIA_DIR / "img-000.png"
PADDLE = MEDIA_DIR / "img-002.png"
SIM_SCREENSHOT = MEDIA_DIR / "img-004.png"
SIM_VIDEO = MEDIA_DIR / "sim_demo.mp4"
VISER_IMG = MEDIA_DIR / "viser_capture.png"


# ----- theme -----------------------------------------------------------------

# Slide is 16:9.
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

NAVY = RGBColor(0x0B, 0x2A, 0x4A)
GREEN = RGBColor(0x4A, 0x7B, 0x5F)
ORANGE = RGBColor(0xD0, 0x80, 0x3C)
TEXT = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x55, 0x55, 0x55)
LIGHT = RGBColor(0xEE, 0xEE, 0xEE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DIVIDER = RGBColor(0xCC, 0xCC, 0xCC)


# ----- helpers ---------------------------------------------------------------

def set_slide_bg(slide, color: RGBColor) -> None:
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill=None, line=None, shape=MSO_SHAPE.RECTANGLE):
    s = slide.shapes.add_shape(shape, x, y, w, h)
    s.shadow.inherit = False
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(0.75)
    s.text_frame.text = ""
    return s


def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=TEXT,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else list(text)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, items, *, size=16, color=TEXT, bullet=True,
                spacing=Pt(6)):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = spacing
        run = p.add_run()
        text = item if isinstance(item, str) else item[0]
        run.text = (f"•  {text}" if bullet else text)
        run.font.name = "Calibri"
        run.font.size = Pt(size)
        run.font.color.rgb = color
        # Optional second-tuple = bold prefix override
        if not isinstance(item, str) and len(item) > 1:
            for tail in item[1:]:
                run2 = p.add_run()
                run2.text = "   " + tail
                run2.font.name = "Calibri"
                run2.font.size = Pt(size - 2)
                run2.font.italic = True
                run2.font.color.rgb = MUTED
    return tb


def add_image(slide, path: Path, x, y, w=None, h=None):
    if not path.exists():
        return None
    kwargs = {}
    if w is not None:
        kwargs["width"] = w
    if h is not None:
        kwargs["height"] = h
    return slide.shapes.add_picture(str(path), x, y, **kwargs)


def add_arrow(slide, x1, y1, x2, y2, color=NAVY, weight=Pt(1.5)):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    c.line.color.rgb = color
    c.line.width = weight
    # Arrowhead on the end of the connector.
    line_elem = c.line._get_or_add_ln()
    from pptx.oxml.ns import qn
    from lxml import etree
    tail = etree.SubElement(line_elem, qn("a:tailEnd"))
    tail.set("type", "triangle")
    tail.set("w", "med")
    tail.set("h", "med")
    return c


def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    # Accent bar
    add_rect(slide, Inches(0), Inches(2.6), SLIDE_W, Inches(0.06), fill=GREEN)
    add_text(slide, Inches(0.9), Inches(1.2), Inches(11.5), Inches(1.4),
             title, size=54, bold=True, color=NAVY)
    add_text(slide, Inches(0.9), Inches(2.85), Inches(11.5), Inches(0.7),
             subtitle, size=24, color=MUTED)
    return slide


def new_content_slide(prs, title, *, eyebrow=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    # Top accent bar
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.18), fill=NAVY)
    # Eyebrow tag
    if eyebrow:
        add_text(slide, Inches(0.6), Inches(0.35), Inches(8), Inches(0.3),
                 eyebrow.upper(), size=11, bold=True, color=GREEN)
        add_text(slide, Inches(0.6), Inches(0.65), Inches(12), Inches(0.7),
                 title, size=30, bold=True, color=NAVY)
    else:
        add_text(slide, Inches(0.6), Inches(0.45), Inches(12), Inches(0.8),
                 title, size=30, bold=True, color=NAVY)
    # Divider
    add_rect(slide, Inches(0.6), Inches(1.30), Inches(12.1), Emu(7000),
             fill=DIVIDER)
    return slide


def add_footer(slide, idx, total):
    add_text(slide, Inches(11.5), Inches(7.15), Inches(1.5), Inches(0.3),
             f"{idx} / {total}", size=10, color=MUTED, align=PP_ALIGN.RIGHT)
    add_text(slide, Inches(0.6), Inches(7.15), Inches(8), Inches(0.3),
             "CS225A — PickleBot — Progress Update", size=10, color=MUTED)


# ----- slide builders --------------------------------------------------------

def slide_title(prs):
    s = add_title_slide(
        prs,
        "PickleBot — Progress Update",
        "CS225A Sports Robotics  ·  Mobile-manipulator pickleball returner",
    )
    # Small subline
    add_text(s, Inches(0.9), Inches(3.6), Inches(11.5), Inches(0.6),
             "Hemal Arora", size=18, color=TEXT)
    add_text(s, Inches(0.9), Inches(4.05), Inches(11.5), Inches(0.5),
             "May 2026", size=14, color=MUTED)
    # Decorative chip
    add_image(s, PADDLE, Inches(10.2), Inches(4.0), h=Inches(2.6))
    return s


def slide_goal(prs):
    s = new_content_slide(prs, "What we're building", eyebrow="Goal & scope")
    add_bullets(s, Inches(0.6), Inches(1.65), Inches(12.0), Inches(5.0), [
        "A mobile-manipulator pickleball returner: TidyBot omni base + Franka Panda arm + custom-mounted MT-01 paddle.",
        "Goal: detect an incoming ball with motion capture, predict where it crosses a strike plane in front of the robot, and swing the paddle through that point to return it.",
        "Today: full sim loop is end-to-end (sim ball  →  state machine  →  controller  →  arm + base  →  paddle hits the ball).",
        "In progress: real-cart bring-up in the SRC Kitchen — OptiTrack streaming live, ball trajectories being recorded and replayed against the production tracker.",
        "This deck: simulation work, control state machine, end-effector design, sensor calibration, and what's next.",
    ], size=18)
    return s


def slide_architecture(prs):
    s = new_content_slide(prs, "System architecture", eyebrow="One contract: Redis")

    # Three column boxes
    col_y = Inches(2.0)
    col_h = Inches(3.6)
    col_w = Inches(3.6)
    gap = Inches(0.45)
    x0 = Inches(0.6)

    def column(x, header, body_lines, fill):
        add_rect(s, x, col_y, col_w, Inches(0.55), fill=fill)
        add_text(s, x, col_y, col_w, Inches(0.55),
                 header, size=16, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(s, x, col_y + Inches(0.55), col_w, col_h - Inches(0.55),
                 fill=LIGHT)
        add_bullets(
            s, x + Inches(0.18), col_y + Inches(0.72),
            col_w - Inches(0.36), col_h - Inches(0.72),
            body_lines, size=13, bullet=False, spacing=Pt(4),
        )

    column(x0, "OptiTrack (SRC Kitchen)", [
        "Motive PC, 8-camera rig",
        "NatNet → Redis bridge",
        "120 Hz pose, unicast UDP",
        "Z-up streaming, world frame",
        "calibrated by 2D Procrustes",
    ], NAVY)

    column(x0 + col_w + gap, "Redis bus", [
        "sai2::optitrack::rigid_body_pos::<id>",
        "sports_bot::cmd::racket::goal_*",
        "sports_bot::cmd::base::goal_pose",
        "sports_bot::state::racket::current_*",
        "sports_bot::fsm::state",
        " ",
        "Same keys for sim & real cart",
    ], GREEN)

    column(x0 + 2 * (col_w + gap), "Robot stack", [
        "Python FSM @ 100 Hz",
        "  predicts intercept, plans swing",
        "C++ hierarchical controller",
        "  base ▶ racket ▶ posture tasks",
        "TidyBot omni driver",
        "  consumes hb1::desired_pose",
        "Franka arm",
        "  joint torques from controller",
    ], NAVY)

    # Connecting arrows
    y_arrow = col_y + Inches(1.8)
    add_arrow(s, x0 + col_w, y_arrow, x0 + col_w + gap, y_arrow)
    add_arrow(s, x0 + 2 * col_w + gap, y_arrow,
              x0 + 2 * col_w + 2 * gap, y_arrow)

    add_text(s, Inches(0.6), Inches(5.9), Inches(12.0), Inches(1.2), (
        "Everything talks through Redis. Bringing up any one piece "
        "(streamer, FSM, controller) in isolation works because the bus "
        "is the only contract — that's how sim and real share the same FSM."
    ), size=14, color=MUTED)
    return s


def slide_sim_demo(prs):
    s = new_content_slide(prs, "Simulation: full FSM loop", eyebrow="Demo")
    # Try video, fall back to screenshot, fall back to placeholder rectangle.
    placeholder_x = Inches(2.7)
    placeholder_y = Inches(1.7)
    placeholder_w = Inches(8.0)
    placeholder_h = Inches(4.5)

    inserted = False
    if SIM_VIDEO.exists():
        try:
            s.shapes.add_movie(
                str(SIM_VIDEO),
                placeholder_x, placeholder_y, placeholder_w, placeholder_h,
                poster_frame_image=str(SIM_SCREENSHOT) if SIM_SCREENSHOT.exists() else None,
            )
            inserted = True
        except Exception:
            inserted = False
    if not inserted and SIM_SCREENSHOT.exists():
        add_image(s, SIM_SCREENSHOT, placeholder_x, placeholder_y,
                  w=placeholder_w)
        inserted = True
    if not inserted:
        add_rect(s, placeholder_x, placeholder_y, placeholder_w, placeholder_h,
                 fill=LIGHT, line=DIVIDER)
        add_text(s, placeholder_x, placeholder_y, placeholder_w, placeholder_h,
                 "[ sim demo video — drop in sports_bot/media/sim_demo.mp4 ]",
                 size=18, color=MUTED, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)

    add_text(s, Inches(0.6), Inches(6.4), Inches(12.0), Inches(0.45),
             "READY → TRACK → APPROACH → SWING → RECOVER, one full hit",
             size=18, bold=True, color=NAVY)
    add_text(s, Inches(0.6), Inches(6.85), Inches(12.0), Inches(0.4),
             "SAI simviz  ·  hierarchical task controller  ·  Python FSM (100 Hz)  ·  ball publishes to the same OptiTrack-shaped Redis key the real cart uses.",
             size=12, color=MUTED)
    return s


def slide_state_machine(prs):
    s = new_content_slide(prs, "Control state machine", eyebrow="Strategy")

    states = ["INIT", "READY", "TRACK", "APPROACH", "SWING", "RECOVER"]
    n = len(states)
    box_w = Inches(1.75)
    box_h = Inches(0.85)
    y_row = Inches(2.4)
    total_w = n * box_w + (n - 1) * Inches(0.20)
    x_start = (SLIDE_W - total_w) / 2
    positions = []
    for i, name in enumerate(states):
        x = x_start + i * (box_w + Inches(0.20))
        positions.append((x, y_row))
        fill = GREEN if name in ("APPROACH", "SWING") else NAVY
        add_rect(s, x, y_row, box_w, box_h, fill=fill,
                 shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        add_text(s, x, y_row, box_w, box_h, name, size=14, bold=True,
                 color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Forward arrows
    for i in range(n - 1):
        x1 = positions[i][0] + box_w
        x2 = positions[i + 1][0]
        y_mid = y_row + box_h / 2
        add_arrow(s, x1, y_mid, x2, y_mid)
    # Loop arrow RECOVER -> READY (down + left + up)
    rec_x, rec_y = positions[-1]
    rd_x, rd_y = positions[1]
    drop = Inches(0.7)
    y_bot = rec_y + box_h + drop
    # down from recover
    c1 = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                rec_x + box_w / 2, rec_y + box_h,
                                rec_x + box_w / 2, y_bot)
    c1.line.color.rgb = MUTED
    c1.line.width = Pt(1.5)
    # across to under ready
    c2 = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                rec_x + box_w / 2, y_bot,
                                rd_x + box_w / 2, y_bot)
    c2.line.color.rgb = MUTED
    c2.line.width = Pt(1.5)
    # up to ready (with arrowhead)
    add_arrow(s, rd_x + box_w / 2, y_bot, rd_x + box_w / 2, rd_y + box_h,
              color=MUTED, weight=Pt(1.5))
    add_text(s, rec_x - Inches(0.8), y_bot - Inches(0.1), Inches(2.5),
             Inches(0.35), "settled OR 1.5 s timeout", size=10, color=MUTED,
             align=PP_ALIGN.CENTER)

    # SAFE_STOP off to the side
    ss_x = Inches(0.7)
    ss_y = Inches(5.7)
    add_rect(s, ss_x, ss_y, Inches(1.75), Inches(0.7),
             fill=ORANGE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(s, ss_x, ss_y, Inches(1.75), Inches(0.7), "SAFE_STOP",
             size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, ss_x + Inches(1.85), ss_y + Inches(0.05),
             Inches(4.5), Inches(0.6),
             "Reached on any unhandled error.\nHolds the ready pose; no base motion.",
             size=11, color=MUTED)

    # Transition labels under each arrow
    edge_labels = [
        "racket at home",
        "ball incoming + intercept valid",
        "swing plan reachable",
        "≤ 0.20 s to impact",
        "follow-through held",
    ]
    for i, label in enumerate(edge_labels):
        x1 = positions[i][0] + box_w
        x2 = positions[i + 1][0]
        add_text(s, x1, y_row + box_h + Inches(0.08), x2 - x1, Inches(0.55),
                 label, size=9, color=MUTED, align=PP_ALIGN.CENTER)

    # Right-side timing box
    info_x = Inches(7.6)
    info_y = Inches(5.55)
    info_w = Inches(5.2)
    info_h = Inches(1.55)
    add_rect(s, info_x, info_y, info_w, info_h, fill=LIGHT)
    add_text(s, info_x + Inches(0.2), info_y + Inches(0.1),
             info_w - Inches(0.4), Inches(0.35),
             "Key timing (config.py)", size=12, bold=True, color=NAVY)
    add_bullets(s, info_x + Inches(0.2), info_y + Inches(0.45),
                info_w - Inches(0.4), info_h - Inches(0.55), [
        "FSM loop: 100 Hz",
        "Swing commit window: 0.20 s before impact",
        "Ball-loss timeout: 0.40 s",
        "Strike plane at x = 0.60 m in front of robot",
    ], size=11, color=TEXT, bullet=True, spacing=Pt(2))
    return s


def slide_tracker(prs):
    s = new_content_slide(prs, "Ball tracker  &  intercept prediction",
                          eyebrow="Perception")
    # Two-column layout
    left_x = Inches(0.6)
    right_x = Inches(7.0)
    col_w = Inches(6.0)

    add_text(s, left_x, Inches(1.55), col_w, Inches(0.4),
             "How", size=15, bold=True, color=GREEN)
    add_bullets(s, left_x, Inches(2.0), col_w, Inches(4.6), [
        "Sliding window of the most recent 12 ball samples (last 0.30 s).",
        "Per-axis ballistic fit: constant velocity in x/y, gravity-corrected linear fit in z.",
        "Reject blatant jumps (>0.5 m between consecutive samples) — kills OptiTrack mislabels.",
        "Propagate the fit forward to the strike plane in front of the robot.",
        "If the predicted arc would cross the floor first, reflect off the floor with a measured restitution and tangential damping before continuing — so we can hit groundstrokes, not just volleys.",
    ], size=14)

    add_text(s, right_x, Inches(1.55), col_w, Inches(0.4),
             "Why it matters", size=15, bold=True, color=GREEN)
    add_bullets(s, right_x, Inches(2.0), col_w, Inches(2.4), [
        "Real pickleball returns are mostly post-bounce groundstrokes.",
        "A single-arc fit silently averages pre- and post-bounce velocities when the window straddles a bounce → garbage predictions for ~3 ticks.",
        "We detect bounces inside the rolling window and drop pre-bounce samples so the next fit only sees the new arc.",
    ], size=13)

    # Results box
    res_x = right_x
    res_y = Inches(4.85)
    res_w = col_w
    res_h = Inches(2.2)
    add_rect(s, res_x, res_y, res_w, res_h, fill=LIGHT)
    add_text(s, res_x + Inches(0.2), res_y + Inches(0.15),
             res_w - Inches(0.4), Inches(0.4),
             "Measured on real recordings (12 sessions, 19 throws, SRC Kitchen)",
             size=12, bold=True, color=NAVY)
    add_bullets(s, res_x + Inches(0.2), res_y + Inches(0.55),
                res_w - Inches(0.4), res_h - Inches(0.6), [
        "Floor model: restitution e ≈ 0.70,  tangential damping μ_t ≈ 0.62  (from 11 bounces).",
        "Commit-window prediction error (the one the FSM acts on):",
        "       mean  8.2 → 5.0 cm   (−40%)",
        "       max   73  → 31  cm   (−58%)",
    ], size=12, bullet=True, spacing=Pt(3))
    return s


def slide_swing(prs):
    s = new_content_slide(prs, "Swing planner",
                          eyebrow="Strategy")
    add_text(s, Inches(0.6), Inches(1.55), Inches(12), Inches(0.5),
             "Predicted intercept + return target  →  three racket poses + a strike velocity.",
             size=15, color=TEXT)

    # Three pose cards
    card_y = Inches(2.3)
    card_h = Inches(2.6)
    card_w = Inches(3.9)
    gap = Inches(0.25)
    x0 = Inches(0.6)

    def card(x, title, body):
        add_rect(s, x, card_y, card_w, Inches(0.5), fill=NAVY)
        add_text(s, x, card_y, card_w, Inches(0.5), title,
                 size=14, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(s, x, card_y + Inches(0.5), card_w, card_h - Inches(0.5),
                 fill=LIGHT)
        add_bullets(s, x + Inches(0.2), card_y + Inches(0.6),
                    card_w - Inches(0.4), card_h - Inches(0.6),
                    body, size=12, bullet=True, spacing=Pt(3))

    card(x0, "Wind-up", [
        "Pull racket back 0.25 m along −face-normal",
        "Same orientation as strike — minimizes rotation during the swing",
        "Reached during APPROACH",
    ])
    card(x0 + (card_w + gap), "Strike", [
        "Position = predicted intercept point",
        "Face normal points toward return target, lifted to clear the net",
        "Velocity = 4 m/s along +face-normal",
    ])
    card(x0 + 2 * (card_w + gap), "Follow-through", [
        "Push 0.25 m past intercept along +face-normal",
        "Held briefly (0.20 s) before RECOVER",
        "Keeps contact long enough to actually impart velocity",
    ])

    # Base placement
    add_text(s, Inches(0.6), Inches(5.2), Inches(12), Inches(0.4),
             "Mobile base", size=15, bold=True, color=GREEN)
    add_bullets(s, Inches(0.6), Inches(5.6), Inches(12), Inches(1.4), [
        "Place base directly behind the strike point, clamped to safe court range  (x ∈ [−0.3, 0.3] m,  y ∈ [−1.5, 1.5] m).",
        "Plan is re-solved on every tick of APPROACH as new ball samples arrive — late corrections are cheap.",
        "Today the strike plane is fixed at x = 0.60 m in front of the base; reachable hitting region will be measured on the real cart and folded into these limits.",
    ], size=13)
    return s


def slide_mount(prs):
    s = new_content_slide(prs, "End-effector: custom paddle mount",
                          eyebrow="Hardware")

    # Left: mount image
    add_image(s, MOUNT_ISO, Inches(0.6), Inches(1.7), h=Inches(4.7))
    add_text(s, Inches(0.6), Inches(6.45), Inches(6.0), Inches(0.35),
             "Flange adapter, isometric", size=11, color=MUTED,
             align=PP_ALIGN.CENTER)

    # Right: paddle image + spec
    add_image(s, PADDLE, Inches(8.6), Inches(1.7), h=Inches(3.6))

    spec_y = Inches(5.5)
    add_text(s, Inches(7.3), spec_y, Inches(5.6), Inches(0.4),
             "Design", size=14, bold=True, color=GREEN)
    add_bullets(s, Inches(7.3), spec_y + Inches(0.35),
                Inches(5.6), Inches(1.7), [
        "Bolts to the Panda flange (link7).",
        "Hex socket grips the paddle handle; cross-pins lock it in.",
        "Paddle face center sits 0.368 m from link7 origin — that's the sweet-spot the controller tracks.",
    ], size=11, bullet=True, spacing=Pt(2))

    return s


def slide_optitrack(prs):
    s = new_content_slide(prs, "OptiTrack calibration",
                          eyebrow="Sensor pipeline")

    col_w = Inches(4.0)
    col_y = Inches(1.65)
    col_h = Inches(4.4)
    gap = Inches(0.25)
    x0 = Inches(0.6)

    def card(x, header, body):
        add_rect(s, x, col_y, col_w, Inches(0.55), fill=NAVY)
        add_text(s, x, col_y, col_w, Inches(0.55), header,
                 size=13, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(s, x, col_y + Inches(0.55), col_w, col_h - Inches(0.55),
                 fill=LIGHT)
        add_bullets(s, x + Inches(0.18), col_y + Inches(0.7),
                    col_w - Inches(0.36), col_h - Inches(0.7),
                    body, size=11, bullet=True, spacing=Pt(3))

    card(x0, "1.  Network", [
        "Motive PC streams NatNet over UDP.",
        "From Stanford-wifi → SRC subnet, multicast doesn't route — switched the server to unicast.",
        "~8 ms round-trip; 120 Hz pose, no drops.",
        "Drop-in: a NatNet → Redis bridge publishes pose + orientation per rigid body.",
    ])

    card(x0 + col_w + gap, "2.  Frames", [
        "Forced Motive's streaming up-axis to Z so it matches our world frame.",
        "Found a silent display-vs-streaming up-axis mismatch in the GUI — added a one-line floor-marker sanity check we run at session start.",
        "World frame: +X toward opponent, +Y left, +Z up, origin at robot home.",
    ])

    card(x0 + 2 * (col_w + gap), "3.  World transform", [
        "Three taped floor markers: origin, +1 m forward, +1 m right.",
        "2D Procrustes solve → single yaw + translation that takes OptiTrack room frame → world frame.",
        "Residuals: max  0.34 mm  horizontal,  1.87 mm  vertical (floor unevenness, not noise).",
        "Stored in world_calibration.json, hot-loaded by the streamer.",
    ])

    add_text(s, Inches(0.6), Inches(6.25), Inches(12), Inches(0.9), (
        "Result: every consumer of OptiTrack reads world-frame positions on a single Redis key  —  same key the sim ball publishes to. "
        "Re-solve per bay or after any camera recalibration."
    ), size=12, color=MUTED)
    return s


def slide_real_tracking(prs):
    s = new_content_slide(prs, "Real-world ball tracking",
                          eyebrow="Replay & validation")

    # Left: viser image or placeholder
    img_x = Inches(0.6)
    img_y = Inches(1.65)
    img_w = Inches(6.4)
    img_h = Inches(4.6)
    if VISER_IMG.exists():
        add_image(s, VISER_IMG, img_x, img_y, w=img_w)
    else:
        add_rect(s, img_x, img_y, img_w, img_h, fill=LIGHT, line=DIVIDER)
        add_text(s, img_x, img_y, img_w, img_h,
                 "[ viser capture — drop in sports_bot/media/viser_capture.png ]",
                 size=14, color=MUTED, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, img_x, img_y + img_h + Inches(0.05), img_w, Inches(0.35),
             "viser replay: ball samples, fit window, predicted arc, predicted intercepts",
             size=10, color=MUTED, align=PP_ALIGN.CENTER)

    # Right: workflow + numbers
    rx = Inches(7.4)
    add_text(s, rx, Inches(1.65), Inches(5.5), Inches(0.4),
             "Recording harness", size=14, bold=True, color=GREEN)
    add_bullets(s, rx, Inches(2.05), Inches(5.5), Inches(2.2), [
        "Record ball pose from Redis → .npz files; we have 12 sessions, 19 throws, 12 k samples.",
        "Replay each recording through the production ball tracker — no copy of the code, the test harness shares the same fit/predict path.",
        "Visualize in viser: tune parameters and compare predicted-vs-actual intercepts without re-throwing.",
    ], size=12)

    add_text(s, rx, Inches(4.55), Inches(5.5), Inches(0.4),
             "What this told us", size=14, bold=True, color=GREEN)
    add_bullets(s, rx, Inches(4.95), Inches(5.5), Inches(2.3), [
        "Single-arc fit fails across bounces — predictions dive below the floor.",
        "Measured floor bounce parameters from the recordings (e, μ_t) and folded them into the predictor.",
        "Online bounce-aware history pruning closed the loop: 8.2 → 5.0 cm mean error at swing-commit.",
    ], size=12)
    return s


def slide_next(prs):
    s = new_content_slide(prs, "Next steps", eyebrow="Plan")

    # Three columns of next steps
    col_w = Inches(4.0)
    col_y = Inches(1.7)
    col_h = Inches(4.6)
    gap = Inches(0.25)
    x0 = Inches(0.6)

    def card(x, header, body, accent=NAVY):
        add_rect(s, x, col_y, col_w, Inches(0.55), fill=accent)
        add_text(s, x, col_y, col_w, Inches(0.55), header,
                 size=14, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(s, x, col_y + Inches(0.55), col_w, col_h - Inches(0.55),
                 fill=LIGHT)
        add_bullets(s, x + Inches(0.2), col_y + Inches(0.75),
                    col_w - Inches(0.4), col_h - Inches(0.75),
                    body, size=12, bullet=True, spacing=Pt(4))

    card(x0, "Probabilistic ball estimator", [
        "Move from a sliding-window least-squares fit to a Kalman filter over ball position + velocity.",
        "Treat each floor bounce as a discrete state transition that reflects v_z, scales v_xy, and inflates the covariance using the bounce-parameter variance we measured.",
        "Output: a predicted intercept with an uncertainty ellipsoid that shrinks with samples and jumps at each bounce.",
        "Lets the FSM commit to SWING when the prediction is tight enough, instead of on a fixed time-to-impact threshold.",
    ], accent=GREEN)

    card(x0 + col_w + gap, "Real cart, supervised", [
        "Move the robot to commanded poses in the SRC Kitchen with workspace safeguards — soft virtual walls around the cart's footprint so it can't drift into kitchen walls or bystanders.",
        "Measure the actual reachable hitting region on the cart (joint limits + base clearance + paddle reach) and feed it back into the swing planner's strike-point bounds.",
        "End-to-end: stream real OptiTrack ball into the existing FSM, log everything.",
    ])

    card(x0 + 2 * (col_w + gap), "Loose ends", [
        "Bridge the FSM's base goal directly to the TidyBot driver's redis key.",
        "Per-session calibration sanity check against a fixed floor reference rigid body.",
        "Standardize the pickleball's OptiTrack streaming ID across all rigs.",
        "Tune controller gains on the real arm with the paddle mounted.",
    ])
    return s


def slide_thanks(prs):
    s = new_content_slide(prs, "Questions?", eyebrow="Thanks")
    add_text(s, Inches(0.6), Inches(2.4), Inches(12.0), Inches(0.7),
             "Happy to dig into anything in more depth.",
             size=24, color=TEXT)
    add_text(s, Inches(0.6), Inches(3.3), Inches(12.0), Inches(2.0), (
        "Most useful comments for us right now:\n"
        "  •  Recommendations for safely commissioning the cart in the kitchen.\n"
        "  •  Whether the Kalman-with-bounce-jumps direction is the right next step,\n"
        "     or whether we should keep iterating on the least-squares predictor."
    ), size=18, color=MUTED)
    return s


# ----- main ------------------------------------------------------------------

def build() -> Path:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slides = [
        slide_title,
        slide_goal,
        slide_architecture,
        slide_sim_demo,
        slide_state_machine,
        slide_tracker,
        slide_swing,
        slide_mount,
        slide_optitrack,
        slide_real_tracking,
        slide_next,
        slide_thanks,
    ]
    total = len(slides)
    for i, fn in enumerate(slides, 1):
        s = fn(prs)
        if i > 1:
            add_footer(s, i, total)

    prs.save(str(OUTPUT))
    return OUTPUT


if __name__ == "__main__":
    out = build()
    print(f"wrote: {out}")
